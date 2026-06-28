import fitz
from docx import Document
from pptx import Presentation
import pandas as pd


def extract_pdf(path):

    pdf = fitz.open(path)

    text = ""

    for page_no, page in enumerate(pdf):

        text += f"\n\n========== PAGE {page_no+1} ==========\n\n"

        text += page.get_text("text")

    pdf.close()

    return text


# TXT
def extract_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


# DOCX
def extract_docx(path):
    doc = Document(path)

    text = []

    for para in doc.paragraphs:
        text.append(para.text)

    return "\n".join(text)


# PPTX
def extract_pptx(path):
    prs = Presentation(path)

    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


# CSV
def extract_csv(path):
    df = pd.read_csv(path)

    return df.to_string()


# XLSX
def extract_xlsx(path):
    df = pd.read_excel(path)

    return df.to_string()


# Main Function
def extract_text(path):

    extension = path.split(".")[-1].lower()

    if extension == "pdf":
        return extract_pdf(path)

    elif extension == "txt":
        return extract_txt(path)

    elif extension == "docx":
        return extract_docx(path)

    elif extension == "pptx":
        return extract_pptx(path)

    elif extension == "csv":
        return extract_csv(path)

    elif extension == "xlsx":
        return extract_xlsx(path)

    else:
        raise ValueError(
            f"Unsupported file type: {extension}"
        )